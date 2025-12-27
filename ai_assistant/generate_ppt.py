from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Helper to create a slide with Title and Bullet points ---
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Bullet layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.shapes.placeholders[1].text_frame
        
        for ui, text in enumerate(content_text_list):
            if ui == 0:
                p = tf.text = text
            else:
                p = tf.add_paragraph()
                p.text = text

    # --- Slide 1: Title ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "Personal AI Assistant V2"
    subtitle.text = "Hackathon Edition: Vision, Web & Autonomy\n\nTeam Omnios"

    # --- Slide 2: The Problem ---
    add_slide("The Problem", [
        "❌ Context Switching Tax: Developers waste time switching between terminal, browser, and IDE.",
        "❌ Disconnected Tools: Terminal outputs don't understand screen context.",
        "❌ Static CLIs: Traditional command lines are text-only and reactive, not proactive."
    ])

    # --- Slide 3: The Solution ---
    add_slide("The Solution: Agentic AI", [
        "✅ Multimodal Intelligence: An assistant that SEES what you see (Omni-Vision).",
        "✅ Connected Knowledge: Real-time Web Search integration (DuckDuckGo).",
        "✅ System Authority: Direct control over files, apps, and computer states.",
        "✅ Voice & Text: Seamless interaction via speech or keyboard."
    ])

    # --- Slide 4: Key Features ---
    add_slide("Key Features", [
        "👁️ Omni-Vision: 'Explain this error' -> Captures screen & Diagnoses.",
        "🌐 Web Agent: 'Who won the match?' -> Fetches live data.",
        "🛡️ Safety Sandbox: All file operations are restricted to a workspace.",
        "✨ Premium UI: Built with 'Rich' for a futuristic terminal aesthetic."
    ])

    # --- Slide 5: Tech Stack ---
    add_slide("Tech Stack", [
        "🐍 Python 3.11: The core engine.",
        "🧠 OpenAI GPT-4o: The brain (Logic & Vision).",
        "🎨 Rich: TUI (Text User Interface) rendering.",
        "🖱️ PyAutoGUI: Screen capture and control.",
        "🦆 DuckDuckGo Search: Web retrieval."
    ])

    # --- Slide 6: Future Roadmap ---
    add_slide("Future Roadmap", [
        "🚀 Autonomous coding (Self-healing agents).",
        "🚀 Local LLM support (Llama 3 / Mistral).",
        "🚀 Full Desktop GUI automation (Semantic Action policies)."
    ])

    prs.save('Hackathon_Pitch.pptx')
    print("Successfully created 'Hackathon_Pitch.pptx'")

if __name__ == "__main__":
    create_presentation()
